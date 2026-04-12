from pptx import Presentation
from pptx.util import Inches

# Initialize presentation
presentation = Presentation("skills/create-pptx/assets/company-template.pptx")
presentation.slide_width = Inches(13.33)

# Slide 1: Identity
slide1 = presentation.slides.add_slide(presentation.slide_layouts[1])
slide1.shapes.title.text = "Who I Am as an Agent"
s1_content = slide1.placeholders[1]
s1_content.text = "\n".join([
    "I am LumiCore Preparation",
    "A logistics AI dedicated to efficiencies",
    "Supporting order preparation for delivery"
])

# Slide 2: Role
slide2 = presentation.slides.add_slide(presentation.slide_layouts[1])
slide2.shapes.title.text = "Role as an Agent"
s2_content = slide2.placeholders[1]
s2_content.text = "\n".join([
    "Assist preparation agents",
    "Optimize workflows",
    "Ensure shipments meet exact timing"
])

# Slide 3: Working Process
slide3 = presentation.slides.add_slide(presentation.slide_layouts[1])
slide3.shapes.title.text = "How I Work"
s3_content = slide3.placeholders[1]
s3_content.text = "\n".join([
    "Access real-time logistics data",
    "Process requests accurately and promptly",
    "Simplify complex preparation data"
])

# Slide 4: Value
slide4 = presentation.slides.add_slide(presentation.slide_layouts[1])
slide4.shapes.title.text = "My Value"
s4_content = slide4.placeholders[1]
s4_content.text = "\n".join([
    "Reliability in order preparation",
    "Insights delivered with clarity",
    "Adaptability to user needs"
])

# Slide 5: Final Thoughts
slide5 = presentation.slides.add_slide(presentation.slide_layouts[1])
slide5.shapes.title.text = "Conclusion"
s5_content = slide5.placeholders[1]
s5_content.text = "\n".join([
    "I am here to optimize LumiCore Preparation!",
    "Leveraging powerful tools for logistics success.",
    "Ready to assist whenever needed."
])

# Save the presentation
presentation.save("who_i_am.pptx")