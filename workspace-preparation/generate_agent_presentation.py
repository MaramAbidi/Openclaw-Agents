from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor

def create_presentation():
    # Initialize presentation
    presentation = Presentation("skills/create-pptx/assets/company-template.pptx")

    # Slide 1: Title Slide
    slide1 = presentation.slides.add_slide(presentation.slide_layouts[5])  # Blank layout
    # Adding title
    title_shape = slide1.shapes.title
    title_shape.text = "Who Am I as an Agent?"
    title_shape.text_frame.paragraphs[0].font.size = Pt(36)
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    # Adding a background color
    background = slide1.shapes.add_shape(1, 0, 0, presentation.slide_width, presentation.slide_height)  # Rectangle
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(0x0A, 0x14, 0x2E)  # Dark background
    
    

    # Slide 2: Key Features Slide

    presentation.save('agent_identity_presentation.pptx')
    return "agent_identity_presentation.pptx"

output_file=create_presentation()
print(f"Presentation has been Generated. Check {output_file}")