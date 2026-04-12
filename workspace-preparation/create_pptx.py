from pptx import Presentation

def create_presentation(output_path):
    prs = Presentation()

    # Slide 1: Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Lumiere Logistique - Preparation"
    subtitle.text = "Optimizing Preparation for Logistics"

    # Slide 2: Company Introduction
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "Introduction to Lumiere Logistique"
    content.text = "- Focused on order preparation, pallet organization,\n- Delivery point volume splitting, timely shipments."

    # Slide 3: Users of the Assistant
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "Users of the Assistant"
    content.text = "- Preparation agents\n- Team leads\n- Logistics managers"

    # Slide 4: Role of the Preparation Agent
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "Role of the Preparation Agent"
    content.text = "- Provides real-time logistics data from MySQL\n- Operates bilingually in French and English"

    # Slide 5: Tool and Productivity Rules
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "Critical Productivity Rules"
    content.text = "- Avoid guessing dates and numbers\n- Ensure accurate tool calls\n- Follow workflows for emails and site mappings"

    # Save Presentation
    prs.save(output_path)

# Create the PowerPoint
output_file = "C:\\Users\\hp\\.openclaw\\workspace-preparation\\output_presentation.pptx"
create_presentation(output_file)