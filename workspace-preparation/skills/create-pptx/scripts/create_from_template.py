from __future__ import annotations

import argparse
import re
from pathlib import Path
from typing import List, Sequence, Tuple

from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.enum.text import MSO_AUTO_SIZE


SlideSpec = Tuple[str, List[str]]


def parse_content_file(content_path: Path) -> List[SlideSpec]:
    """
    Parse a simple text or markdown content source into slides.

    Format:
    - Blocks are separated by one or more blank lines.
    - In each block, first line is slide title.
    - Remaining lines are bullet points.
    - Leading '-', '*', or numbered prefixes are removed from bullet lines.
    """
    raw = content_path.read_text(encoding="utf-8")
    if content_path.suffix.lower() in {".md", ".markdown"}:
        slides = parse_markdown_sections(raw)
    else:
        slides = parse_block_sections(raw)

    if not slides:
        raise ValueError(f"No slide content found in {content_path}")

    return slides


def parse_block_sections(raw: str) -> List[SlideSpec]:
    lines = [line.rstrip() for line in raw.splitlines()]

    blocks: List[List[str]] = []
    current: List[str] = []

    for line in lines:
        if line.strip():
            current.append(line)
        elif current:
            blocks.append(current)
            current = []

    if current:
        blocks.append(current)

    slides: List[SlideSpec] = []
    for block in blocks:
        title = clean_text(block[0].strip().lstrip("#").strip(), for_title=True)
        bullets: List[str] = []
        for item in block[1:]:
            clean = item.strip()
            if not clean:
                continue
            clean = strip_bullet_prefix(clean)
            clean = clean_text(clean)
            if clean:
                bullets.append(clean)
        slides.append((title, bullets))

    return slides


def parse_markdown_sections(raw: str) -> List[SlideSpec]:
    lines = raw.splitlines()
    slides: List[SlideSpec] = []
    current_title: str | None = None
    current_bullets: List[str] = []
    in_code = False
    in_frontmatter = False

    def flush_current() -> None:
        nonlocal current_title, current_bullets
        if current_title:
            slides.append((clean_text(current_title, for_title=True), current_bullets[:]))
        current_title = None
        current_bullets = []

    for index, line in enumerate(lines):
        stripped = line.strip()

        if index == 0 and stripped == "---":
            in_frontmatter = True
            continue
        if in_frontmatter:
            if stripped == "---":
                in_frontmatter = False
            continue

        if stripped.startswith("```"):
            in_code = not in_code
            continue
        if in_code:
            continue

        heading_match = re.match(r"^(#{1,3})\s+(.+)$", stripped)
        if heading_match:
            flush_current()
            current_title = heading_match.group(2)
            continue

        if not stripped:
            continue

        if current_title is None:
            current_title = "Overview"

        bullet_match = re.match(r"^(?:[-*+]\s+|\d+[.)]\s+)(.+)$", stripped)
        if bullet_match:
            bullet_text = clean_text(bullet_match.group(1))
        else:
            bullet_text = clean_text(stripped)

        if bullet_text:
            current_bullets.append(bullet_text)

    flush_current()

    if slides:
        return slides

    # Fallback for non-structured markdown.
    return parse_block_sections(raw)

def strip_bullet_prefix(value: str) -> str:
    if value.startswith(("- ", "* ")):
        return value[2:].strip()

    # Handle numbered bullets like "1. text"
    dot_index = value.find(".")
    if dot_index > 0 and value[:dot_index].isdigit():
        return value[dot_index + 1 :].strip()

    return value


def clean_text(value: str, for_title: bool = False) -> str:
    text = value.strip()
    if not text:
        return ""

    # Convert markdown links: [label](url) -> label
    text = re.sub(r"\[([^\]]+)\]\([^\)]+\)", r"\1", text)
    # Remove inline code ticks
    text = text.replace("`", "")

    # Replace long snake_case-like tokens with spaced words for readability.
    words: List[str] = []
    for token in text.split():
        if "_" in token and len(token) >= 16:
            token = token.replace("_", " ")
        words.append(token)
    text = " ".join(words)

    text = re.sub(r"\s+", " ", text).strip(" -:")
    if not text:
        return ""

    limit = 90 if for_title else 180
    if len(text) > limit:
        text = text[: limit - 1].rstrip() + "…"
    return text


def chunk_slide_specs(
    slides: Sequence[SlideSpec],
    max_bullets_per_slide: int,
    max_slides: int,
) -> List[SlideSpec]:
    out: List[SlideSpec] = []

    for title, bullets in slides:
        if len(out) >= max_slides:
            break

        if not bullets:
            out.append((title, []))
            continue

        step = max(1, max_bullets_per_slide)
        for i in range(0, len(bullets), step):
            if len(out) >= max_slides:
                break
            chunk = bullets[i : i + step]
            chunk_title = title if i == 0 else f"{title} (cont.)"
            out.append((chunk_title, chunk))

    return out


def choose_content_layout(prs: Presentation):
    """
    Prefer a layout with title + body placeholder. Fall back to first layout.
    """
    # Best choice: layout that has both title and body placeholders.
    for layout in prs.slide_layouts:
        has_title = False
        has_body = False
        for placeholder in layout.placeholders:
            try:
                p_type = placeholder.placeholder_format.type
            except Exception:
                continue
            if p_type in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE):
                has_title = True
            if p_type in (PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT):
                has_body = True
        if has_title and has_body:
            return layout

    # Fallback: any layout with at least two placeholders.
    for layout in prs.slide_layouts:
        if len(layout.placeholders) >= 2:
            return layout

    return prs.slide_layouts[0]


def add_content_slides(prs: Presentation, slides: Sequence[SlideSpec]) -> None:
    layout = choose_content_layout(prs)

    for title, bullets in slides:
        slide = prs.slides.add_slide(layout)

        title_placeholder = slide.shapes.title
        if title_placeholder is not None:
            title_placeholder.text = title
            title_tf = title_placeholder.text_frame
            title_tf.word_wrap = True
            title_tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

        body_placeholder = find_body_placeholder(slide.shapes)
        if body_placeholder is None:
            # If template placeholders are not usable, keep title and skip body.
            continue

        text_frame = body_placeholder.text_frame
        text_frame.clear()
        text_frame.word_wrap = True
        text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

        if bullets:
            text_frame.text = bullets[0]
            for bullet in bullets[1:]:
                paragraph = text_frame.add_paragraph()
                paragraph.text = bullet
                paragraph.level = 0


def find_body_placeholder(shapes):
    # Prefer non-title placeholder with text frame.
    for shape in shapes:
        if not getattr(shape, "is_placeholder", False):
            continue
        if shape == shapes.title:
            continue
        if getattr(shape, "has_text_frame", False):
            return shape

    # Fall back to any text box-like shape with a text frame.
    for shape in shapes:
        if getattr(shape, "has_text_frame", False) and shape != shapes.title:
            return shape

    return None


def should_keep_existing_slide(index: int, total: int, keep_cover: bool, keep_last: bool) -> bool:
    if keep_cover and index == 0:
        return True
    if keep_last and index == (total - 1):
        return True
    return False


def remove_unneeded_existing_slides(prs: Presentation, keep_cover: bool, keep_last: bool) -> None:
    # Remove slides from back to front for stable indices.
    total = len(prs.slides)
    for index in reversed(range(total)):
        if should_keep_existing_slide(index, total, keep_cover, keep_last):
            continue
        slide_id = prs.slides._sldIdLst[index]
        rel_id = slide_id.rId
        prs.part.drop_rel(rel_id)
        del prs.slides._sldIdLst[index]


def resolve_path(path_value: str, base_dir: Path) -> Path:
    path = Path(path_value)
    if not path.is_absolute():
        path = base_dir / path
    return path.resolve()


def build_parser() -> argparse.ArgumentParser:
    parser = argparse.ArgumentParser(
        description="Create a PPTX from a template and text/markdown content.")
    parser.add_argument(
        "--template",
        default="../assets/template.pptx",
        help="Path to .pptx template (relative to scripts/ by default).",
    )
    parser.add_argument(
        "--content",
        required=True,
        help="Path to text/markdown content source.",
    )
    parser.add_argument(
        "--output",
        default="../output/generated-from-template.pptx",
        help="Output PPTX path.",
    )
    parser.add_argument(
        "--keep-cover",
        action="store_true",
        help="Keep the first slide from template (often branding cover).",
    )
    parser.add_argument(
        "--keep-last",
        action="store_true",
        help="Keep the last slide from template (often credits/branding).",
    )
    parser.add_argument(
        "--max-bullets",
        type=int,
        default=5,
        help="Maximum bullet points per slide before continuation slide is created.",
    )
    parser.add_argument(
        "--max-slides",
        type=int,
        default=20,
        help="Maximum number of generated content slides.",
    )
    return parser


def main() -> None:
    parser = build_parser()
    args = parser.parse_args()

    scripts_dir = Path(__file__).resolve().parent
    template_path = resolve_path(args.template, scripts_dir)
    content_path = resolve_path(args.content, scripts_dir)
    output_path = resolve_path(args.output, scripts_dir)

    if not template_path.exists():
        raise FileNotFoundError(f"Template not found: {template_path}")
    if not content_path.exists():
        raise FileNotFoundError(f"Content file not found: {content_path}")

    slide_specs = parse_content_file(content_path)
    slide_specs = chunk_slide_specs(
        slide_specs,
        max_bullets_per_slide=max(1, args.max_bullets),
        max_slides=max(1, args.max_slides),
    )

    prs = Presentation(str(template_path))
    remove_unneeded_existing_slides(
        prs,
        keep_cover=args.keep_cover,
        keep_last=args.keep_last,
    )
    add_content_slides(prs, slide_specs)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    print(f"Saved presentation: {output_path}")


if __name__ == "__main__":
    main()
